# app.py (管理者ページ)

import streamlit as st
import sqlite3
import json
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib as mpl
from openai import OpenAI
import numpy as np
from pypdf import PdfReader
from pptx import Presentation
import io
import japanize_matplotlib 


# --- パスワード認証機能 ---
def check_password():
    """Returns `True` if the user had the correct password."""

    # st.secretsからパスワードを読み込む
    # secrets.toml に ADMIN_PASSWORD が設定されていない場合は、この機能は動作しない
    if "ADMIN_PASSWORD" not in st.secrets:
        st.error("パスワードが設定されていません。管理者にお問い合わせください。")
        return False

    # セッション状態で認証済みかどうかをチェック
    if "password_correct" in st.session_state and st.session_state["password_correct"]:
        return True

    # ログインフォームを表示
    st.header("管理者ログイン")
    password = st.text_input("パスワードを入力してください", type="password")
    if st.button("ログイン"):
        if password == st.secrets["ADMIN_PASSWORD"]:
            # パスワードが正しければ、セッション状態を更新して再実行
            st.session_state["password_correct"] = True
            st.rerun()
        else:
            st.error("パスワードが間違っています。")
    return False

# --- メインのアプリケーションロジック ---
try:
    # secrets.tomlからAPIキーを読み込む
    client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])
except Exception:
    st.error("OpenAIのAPIキーがsecrets.tomlに正しく設定されていません。")
    # check_password()の外なので、st.stop()は使わない

# パスワードチェックを実行
if check_password():
    # --- ここから下は、認証が成功した場合にのみ表示・実行される ---

# --- データベースのセットアップ ---
    def init_db():
        conn = sqlite3.connect('qa.db')
        cursor = conn.cursor()
        # questionsテーブルの完全な定義
        cursor.execute('''
        CREATE TABLE IF NOT EXISTS questions (
            id INTEGER PRIMARY KEY AUTOINCREMENT, question_type TEXT NOT NULL, 
            question_text TEXT NOT NULL, options TEXT, answers TEXT NOT NULL, 
            explanation TEXT, difficulty TEXT, points INTEGER, 
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
        ''')
        # answers_historyテーブルの完全な定義
        cursor.execute('''
        CREATE TABLE IF NOT EXISTS answers_history (
            id INTEGER PRIMARY KEY AUTOINCREMENT, user_id TEXT NOT NULL, 
            question_id INTEGER NOT NULL, user_answer TEXT NOT NULL, 
            is_correct BOOLEAN NOT NULL, awarded_score INTEGER, feedback TEXT,
            answered_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (question_id) REFERENCES questions (id)
        )
        ''')
        # ★新規: フィードバックを保存するテーブルを追加
        cursor.execute('''
        CREATE TABLE IF NOT EXISTS test_feedback (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id TEXT NOT NULL,
            satisfaction_rating INTEGER,
            free_text_comment TEXT,
            submitted_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
        ''')
        conn.commit()
        conn.close()

    def extract_text_from_file(uploaded_file):
        try:
            if uploaded_file.name.endswith('.pdf'):
                # PDFファイルの場合
                pdf_reader = PdfReader(uploaded_file)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
                return text
            elif uploaded_file.name.endswith('.pptx'):
                # PowerPointファイルの場合
                prs = Presentation(uploaded_file)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text
            else:
                st.error("対応していないファイル形式です。.pdfまたは.pptxファイルをアップロードしてください。")
                return None
        except Exception as e:
            st.error(f"ファイル読み込み中にエラーが発生しました: {e}")
            return None

    def get_question_stats(user_id=None):
        conn = sqlite3.connect('qa.db'); cursor = conn.cursor()
        query = """
            WITH FirstAttempts AS (
                SELECT
                    h.question_id,
                    h.user_id,
                    h.is_correct,
                    h.awarded_score,
                    q.points,
                    ROW_NUMBER() OVER(PARTITION BY h.user_id, h.question_id ORDER BY h.answered_at ASC) as attempt_rank
                FROM 
                    answers_history h
                JOIN 
                    questions q ON h.question_id = q.id
            )
            SELECT 
                question_id,
                COUNT(*) as attempts,
                SUM(is_correct) as corrects,
                SUM(awarded_score) as total_awarded,
                SUM(points) as total_possible
            FROM 
                FirstAttempts
            WHERE 
                attempt_rank = 1
        """
        params = []
        if user_id:
            query += " AND user_id = ?"
            params.append(user_id)
        query += " GROUP BY question_id"
        
        cursor.execute(query, params)
        stats = cursor.fetchall()
        conn.close()
        stats_dict = {
            row[0]: {
                'attempts': row[1], 
                'corrects': row[2],
                'total_awarded': row[3] if row[3] is not None else 0,
                'total_possible': row[4] if row[4] is not None else 0
            } for row in stats
        }
        return stats_dict
    
    # ★新規: 解答履歴のある全ユーザーを取得する関数
    def get_all_user_ids():
        conn = sqlite3.connect('qa.db'); cursor = conn.cursor()
        cursor.execute("SELECT DISTINCT user_id FROM answers_history ORDER BY user_id")
        users = [row[0] for row in cursor.fetchall()]; conn.close()
        return users

    # ★新規: 全生徒の合計点をリストで取得する関数
    def get_all_student_total_scores():
        conn = sqlite3.connect('qa.db'); cursor = conn.cursor()
        # ユーザーごとに獲得点の合計を計算
        cursor.execute("""
            SELECT SUM(awarded_score)
            FROM answers_history
            GROUP BY user_id
        """)
        scores = [row[0] for row in cursor.fetchall()]
        conn.close()
        return scores
    
    # ★新規：問題を更新する関数
    def update_question(q_id, q_text, q_options, q_answers, q_explanation, q_difficulty, q_points):
        conn = sqlite3.connect('qa.db'); cursor = conn.cursor()
        cursor.execute("UPDATE questions SET question_text=?, options=?, answers=?, explanation=?, difficulty=?, points=? WHERE id=?", (q_text, q_options, q_answers, q_explanation, q_difficulty, q_points, q_id))
        conn.commit(); conn.close()

    # ★新規：問題を削除する関数
    def delete_question(q_id):
        conn = sqlite3.connect('qa.db'); cursor = conn.cursor()
        cursor.execute("DELETE FROM answers_history WHERE question_id=?", (q_id,))
        cursor.execute("DELETE FROM questions WHERE id=?", (q_id,))
        conn.commit(); conn.close()

    # ★新規: 講義フィードバックのデータを取得する関数
    def get_feedback_data():
        conn = sqlite3.connect('qa.db')
        # pd.read_sql_queryを使うと、直接DataFrameとしてデータを読み込めて便利
        query = "SELECT user_id, satisfaction_rating, free_text_comment, submitted_at FROM test_feedback ORDER BY submitted_at DESC"
        df = pd.read_sql_query(query, conn)
        conn.close()
        return df

    # ★新規: 数値の評価を星の表示に変換するヘルパー関数
    def rating_to_stars(rating):
        if not rating or pd.isna(rating):
            return "評価なし"
        # 小数点以下を四捨五入して星の数を決定
        rounded_rating = int(round(rating, 0))
        return "⭐" * rounded_rating

    # ★新規: 偏差値を計算する関数
    def calculate_hensachi(score, mean, std_dev):
        # 標準偏差が0の場合（全員が同点、または受験者が1人など）は偏差値を50とする
        if std_dev == 0:
            return 50.0
        return 50 + 10 * (score - mean) / std_dev

    # ★新規: 統計データからグラフを生成する関数
    def create_performance_chart(individual_stats, overall_stats):
        if not overall_stats: return None
        overall_df = pd.DataFrame.from_dict(overall_stats, orient='index')
        overall_df['overall_score_rate'] = (overall_df['total_awarded'] / overall_df['total_possible']) * 100
        if individual_stats:
            ind_df = pd.DataFrame.from_dict(individual_stats, orient='index')
            ind_df['individual_score_rate'] = (ind_df['total_awarded'] / ind_df['total_possible']) * 100
            # それぞれ必要な「得点率」の列だけを抽出してから結合する
            df_overall_rate = overall_df[['overall_score_rate']]
            df_individual_rate = ind_df[['individual_score_rate']]
            
            df = df_overall_rate.join(df_individual_rate, how='outer').fillna(0)
            chart_data = df[['individual_score_rate', 'overall_score_rate']]; title = "個人 vs 全体平均 得点率比較"
        else:
            df = overall_df; chart_data = df[['overall_score_rate']]; title = "全体平均 得点率"
        ax = chart_data.plot(kind='barh', figsize=(10, len(df) * 0.6 if len(df) > 0 else 3), color=['tomato', 'skyblue'] if individual_stats else 'skyblue', width=0.8)
        # ★追加: グラフの縦軸の順番を逆にします (ID 1が上に来る)
        ax.invert_yaxis()
        ax.set_xlabel('得点率 (%)'); ax.set_ylabel('問題ID'); ax.set_title(title); ax.set_xlim(0, 105); ax.legend(title='凡例')
        for container in ax.containers: ax.bar_label(container, fmt='%.1f%%', padding=3)
        plt.tight_layout(); chart_path = "performance_chart.png"; plt.savefig(chart_path); plt.close(ax.get_figure())
        return chart_path

    # --- QA生成のロジック（変更なし） ---
    def create_qa_set(lecture_text, difficulty_level, num_questions_mode, num_questions_manual=5):
        # モードに応じてプロンプトの「問題数」に関する部分を組み立てる
        if num_questions_mode == "AIにおまかせ":
            count_instruction = "まず、このテキスト全体の要点を把握し、学習者が理解度を確認するのに適切だと思われる問題数を3問から10問の範囲で判断してください。その後、あなたが判断したその問題数で、以下の指示に従って問題セットを生成してください。"
        else: # 手動指定の場合
            count_instruction = f"問題を合計{num_questions_manual}問作成してください。"

        # 共通のプロンプト部分
        prompt = f"""
        # 指示
        - 以下のテキストに基づき、内容の理解度を確認するための問題を作成せよ。
        - {count_instruction}
        - 問題形式は「4択問題」「5択複数選択問題」「自由記述問題」をバランス良く混ぜること。
        - 選択式問題では、問題文に「一つ選びなさい」「二つ選びなさい」「全て選びなさい」といった解答数を指定する文言を必ず含めること。
        - 全体の難易度は「{difficulty_level}」レベルを意識すること。
        - 各問題には "easy", "normal", "hard" のいずれかの難易度を必ず設定すること。
        # 出力形式
        以下の厳格なJSON形式で、結果のみを出力せよ。
        {{
        "questions": [
            {{
            "difficulty": "...",
            "question_type": "...",
            "question_text": "...",
            "options": ["...", "..."],
            "answers": ["...", "..."],
            "explanation": "..."
            }}
        ]
        }}

        # テキスト
        {lecture_text}
        """
    
        try:
            response = client.chat.completions.create(
                model="gpt-4o", messages=[{"role": "user", "content": prompt}],
                response_format={"type": "json_object"}
            )
            qa_data = json.loads(response.choices[0].message.content)
            return qa_data.get("questions", [])
        except Exception as e:
            st.error(f"エラーが発生しました: {e}")
            return []
        
    # ★新規: 全てのデータを統合し、AIに評価と改善策を提案させる関数
    def generate_ai_analysis(scope, score_stats_df, feedback_df, questions_df):
        # AIに渡すためのテキストデータを作成
        prompt_data = f"分析対象: {scope}\n\n"
        
        # 1. 得点率データ
        prompt_data += "--- 問題別の得点率データ ---\n"
        if not score_stats_df.empty:
            for q_id, row in score_stats_df.iterrows():
                question_text = questions_df.loc[q_id, 'question_text']
                score_rate = row.get('score_rate', 0)
                prompt_data += f"- 問題ID {q_id} (内容: {question_text[:30]}...): 得点率 {score_rate:.1f}%\n"
        else:
            prompt_data += "得点データなし\n"
            
        # 2. 満足度・フィードバックデータ
        prompt_data += "\n--- 講義内容へのフィードバック ---\n"
        if not feedback_df.empty:
            avg_rating = feedback_df['satisfaction_rating'].mean()
            prompt_data += f"平均満足度: {avg_rating:.2f} / 5.0\n"
            prompt_data += "自由記述コメント:\n"
            for index, row in feedback_df.iterrows():
                prompt_data += f"- {row['free_text_comment']}\n"
        else:
            prompt_data += "フィードバックデータなし\n"
            
        # AIへの指示（プロンプト）
        final_prompt = f"""
        あなたは優秀な教育データアナリストです。
        講義内容や説明方法の改善策の提案をしてもらいたいです。
        また生徒の正確な理解度把握のため、QAの質やアンケートの取り方についても改善策を提案してもらいたいです。
        以下の全体もしくは個人の生徒の成績データとフィードバックを総合的に分析し、講義内容や説明方式、あるいは生成されたQAの質や実施したアンケートに関する「評価」と「具体的な改善策の提案」を日本語で記述してください。

        # 分析データ
        {prompt_data}

        # 出力形式
        以下のMarkdown形式で、簡潔かつ分かりやすくまとめてください。

        ### 総合評価
        （ここに参加者全体の、あるいは個人の理解度や満足度に関する客観的な評価を記述してください。特に得点率が低い問題や、具体的なフィードバック内容に言及してください。）

        ### 講義内容に関する改善点の提案
        1. （評価に基づいて、講義内容、説明方法の改善に関する具体的なアクションアイテムを提案してください。）
        2. （2つ目の提案）
        3. （3つ目の提案）

        ### QA＆アンケートに関する改善点の提案
        1. （評価に基づいて、生成されたQAや実施したアンケートの改善に関する具体的なアクションアイテムを提案してください。）
        2. （2つ目の提案）
        3. （3つ目の提案）

        """
        
        try:
            client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])
            response = client.chat.completions.create(
                model="gpt-4o",
                messages=[{"role": "user", "content": final_prompt}]
            )
            return response.choices[0].message.content
        except Exception as e:
            return f"AI分析中にエラーが発生しました: {e}"

    # --- StreamlitのUI部分 ---
    init_db()

    if 'editing_question_id' not in st.session_state:
        st.session_state.editing_question_id = None

    st.title("管理者ページ")

    # 2つのタブを作成
    tab_management, tab_analytics = st.tabs(["QA生成・管理", "成績分析"])

    #st.info("ここでは講義テキストから問題を作成し、内容を確認できます。")

    with tab_management:
        st.header("QAの生成と管理")

        with st.expander("新しいQAセットを生成する", expanded=True):
            # タブで入力モードを切り替え
            tab1, tab2 = st.tabs(["テキスト入力で生成", "ファイルアップロードで生成"])

            source_text = None

            with tab1:
                st.subheader("テキストを直接入力")
                text_input = st.text_area("ここに講義テキストを貼り付け", height=250, key="text_input_area")

            with tab2:
                st.subheader("PDFまたはPowerPointファイルをアップロード")
                uploaded_file = st.file_uploader(
                    "講義資料ファイルを選択",
                    type=['pdf', 'pptx'],
                    key="file_uploader_widget"
                )
            
            st.divider()
            st.write("**設定**")
            col1, col2 = st.columns(2)
            with col1:
                # 問題数設定モードの選択
                num_q_mode = st.radio(
                    "問題数の設定",
                    ("手動で指定", "AIにおまかせ"),
                    key="num_q_mode_radio"
                )
            with col2:
                # 手動指定モードの場合のみ、数値入力ボックスを表示
                if num_q_mode == "手動で指定":
                    num_q_manual = st.number_input(
                        "問題数", min_value=1, max_value=20, value=5,
                        key="num_q_manual_input"
                    )
                else:
                    num_q_manual = 5 # 使わないがデフォルト値を設定
                    st.info("AIがテキスト量に応じて問題数を自動で判断します。")
            
            difficulty = st.radio("難易度", ('簡単', '普通', '難しい'), index=1, horizontal=True, key="difficulty_radio")

            if st.button("QAを生成する"):
                # どちらのモードで入力されたか判定
                if uploaded_file is not None:
                    # ファイルモード
                    with st.spinner(f"「{uploaded_file.name}」を読み込んでいます..."):
                        source_text = extract_text_from_file(uploaded_file)
                        if source_text:
                            st.info("ファイルの読み込みが完了しました。")
                elif text_input.strip() != "":
                    # テキストモード
                    source_text = text_input
                else:
                    st.warning("テキストを入力するか、ファイルをアップロードしてください。")

                if source_text:
                    with st.spinner(f"AIが「{difficulty}」レベルのQAセットを生成中です..."):
                        questions_list = create_qa_set(
                            source_text, 
                            difficulty, 
                            num_q_mode, 
                            num_q_manual if num_q_mode == "手動で指定" else None
                        )
                        if questions_list:
                            conn = sqlite3.connect('qa.db'); cursor = conn.cursor()
                            for q in questions_list:
                                # ★追加: 難易度に基づいて配点を決定
                                difficulty = q.get("difficulty", "normal")
                                if difficulty == "easy":
                                    points = 5
                                elif difficulty == "hard":
                                    points = 10
                                else:
                                    points = 7
                            
                                # ★変更: pointsをDBに保存
                                cursor.execute(
                                    "INSERT INTO questions (question_type, question_text, options, answers, explanation, difficulty, points) VALUES (?, ?, ?, ?, ?, ?, ?)",
                                    (q.get("question_type"), q.get("question_text"),
                                    json.dumps(q.get("options")), json.dumps(q.get("answers")),
                                    q.get("explanation"), difficulty, points)
                                )
                            conn.commit(); conn.close()
                            st.success(f"{len(questions_list)}問のQAが生成・保存されました！")
                else:
                    st.warning("講義テキストを入力してください。")

        st.divider()
        st.subheader("📚 全てのQA一覧（管理者ビュー）")

        # セッション状態で編集中の問題IDを管理
        #if 'editing_question_id' not in st.session_state:
            #st.session_state.editing_question_id = None

        # ★追加: 全問題の統計情報を最初に取得
        question_stats = get_question_stats()

        conn = sqlite3.connect('qa.db'); cursor = conn.cursor()
        cursor.execute("SELECT id, question_type, question_text, options, answers, explanation, difficulty, points FROM questions ORDER BY id ASC")
        all_questions = cursor.fetchall(); conn.close()

        if not all_questions:
            st.info("まだ保存されているQAはありません。")
        else:
            for q in all_questions:
                # データベースから取得した各列のデータを、分かりやすい名前の変数に格納します
                q_id, q_type, q_text, q_options_json, q_answers_json, q_explanation, q_difficulty, q_points = q

                stats = question_stats.get(q_id, {'attempts': 0, 'corrects': 0})
                attempts = stats['attempts']
                corrects = stats['corrects']
                # 正答率を計算
                correct_rate = (corrects / attempts * 100) if attempts > 0 else 0

                if st.session_state.editing_question_id == q_id:
                    with st.form(key=f"edit_form_{q_id}"):
                        st.write(f"**ID: {q_id} を編集中...**")
                        
                        # 1. 基本情報の編集
                        new_q_text = st.text_area("問題文", value=q_text)
                        new_q_explanation = st.text_area("解説", value=q_explanation)
                        col1, col2 = st.columns(2)
                        with col1:
                            new_q_difficulty = st.selectbox("難易度", ["easy", "normal", "hard"], index=["easy", "normal", "hard"].index(q_difficulty))
                        with col2:
                            new_q_points = st.number_input("配点", value=q_points)

                        st.write("---")

                        # 2. 選択肢と正解の編集（UI改善）
                        # 選択肢がある問題形式の場合のみ表示
                        if q_type in ["4択問題", "5択複数選択問題"]:
                            st.write("**選択肢と正解の編集**")
                            # DBのJSON文字列をPythonのリストに変換
                            current_options = json.loads(q_options_json) if q_options_json else []
                            # リストを改行区切りのテキストに変換して表示
                            options_as_text = "\n".join(current_options)
                            
                            new_options_text = st.text_area("選択肢 (1行に1つずつ入力)", value=options_as_text, height=150)
                            
                            # 現在入力されている選択肢をリアルタイムでリスト化
                            options_for_multiselect = [opt.strip() for opt in new_options_text.strip().split('\n') if opt.strip()]
                            
                            # DBのJSON文字列から現在の正解リストを取得
                            current_answers = json.loads(q_answers_json) if q_answers_json else []

                            new_answers_list = st.multiselect(
                                "正解を選択 (上記で入力した選択肢から選んでください)",
                                options=options_for_multiselect,
                                default=current_answers
                            )
                        else: # 自由記述問題の場合
                            st.write("**模範解答の編集**")
                            current_answers = json.loads(q_answers_json) if q_answers_json else [""]
                            # 自由記述の答えはリストの最初の要素とする
                            new_single_answer = st.text_area("模範解答", value=current_answers[0])


                        # 3. 保存ボタン
                        submitted = st.form_submit_button("保存する")
                        if submitted:
                            # フォームの入力値をDB保存用の形式に変換
                            final_options_json = None
                            final_answers_json = None
                            if q_type in ["4択問題", "5択複数選択問題"]:
                                # 改行区切りテキストをJSON文字列に戻す
                                final_options_list = [opt.strip() for opt in new_options_text.strip().split('\n') if opt.strip()]
                                final_options_json = json.dumps(final_options_list, ensure_ascii=False)
                                # multiselectの結果をJSON文字列に戻す
                                final_answers_json = json.dumps(new_answers_list, ensure_ascii=False)
                            else: # 自由記述
                                # 答えをリスト形式にしてJSON文字列に戻す
                                final_answers_json = json.dumps([new_single_answer], ensure_ascii=False)

                            # DB更新
                            update_question(q_id, new_q_text, final_options_json, final_answers_json, new_q_explanation, new_q_difficulty, new_q_points)
                            st.session_state.editing_question_id = None
                            st.success(f"ID: {q_id}の問題を更新しました。")
                            st.rerun()
                
                # 通常の表示モード
                else:
                    with st.expander(f"ID: {q_id} | 形式: {q_type} | 難易度: {q_difficulty} | 配点: {q_points}点"):
                        # ★追加: 統計情報を表示
                        col1, col2 = st.columns(2)
                        with col1:
                            st.metric("回答数 (初回挑戦)", f"{attempts} 回")
                        with col2:
                            st.metric("平均正答率 (初回挑戦)", f"{correct_rate:.1f} %")
                        
                        st.divider()
                        
                        st.markdown(f"**問題:** {q_text}")
                        if q_options_json: st.write("**選択肢:**"); st.json(json.loads(q_options_json))
                        st.write("---")
                        if q_answers_json: st.markdown(f"**正解:**"); st.json(json.loads(q_answers_json))
                        st.markdown(f"**解説:** {q_explanation}")

                        # 編集ボタンと削除ボタンを横並びに配置
                        col1, col2 = st.columns(2)
                        with col1:
                            if st.button("編集", key=f"edit_{q_id}"):
                                st.session_state.editing_question_id = q_id
                                st.rerun()
                        with col2:
                            if st.button("削除", key=f"delete_{q_id}"):
                                delete_question(q_id)
                                st.success(f"ID: {q_id}の問題を削除しました。")
                                st.rerun()

    with tab_analytics:
        st.header("成績分析ダッシュボード")

        selected_user = "全体"  # デフォルト値を設定
        user_list = get_all_user_ids()
        selected_user = st.selectbox(
            "分析対象のユーザーを選択",
            ["全体"] + user_list,
            key="analytics_user_selector"  # ユニークなキー
        )

        # 全体平均の統計データを常に取得
        overall_stats_data = get_question_stats()
        all_student_scores = get_all_student_total_scores()

        # 全体の平均と標準偏差を計算
        mean_score = np.mean(all_student_scores) if all_student_scores else 0
        std_dev_score = np.std(all_student_scores) if all_student_scores else 0

        # --- メトリック表示 ---
        col1, col2, col3 = st.columns(3)
        
        with col3:
            total_awarded_overall = sum(d.get('total_awarded', 0) for d in overall_stats_data.values())
            total_possible_overall = sum(d.get('total_possible', 0) for d in overall_stats_data.values())
            overall_rate = (total_awarded_overall / total_possible_overall * 100) if total_possible_overall > 0 else 0
            st.metric(label="全体の平均得点率", value=f"{overall_rate:.1f} %")
        
        individual_stats_data = None
        display_user = "全体"
        if selected_user != "全体":
            individual_stats_data = get_question_stats(selected_user)
        
            # 「個人の平均得点率」を一番左(col1)に配置
            with col1:
                total_awarded_ind = sum(d.get('total_awarded', 0) for d in individual_stats_data.values())
                total_possible_ind = sum(d.get('total_possible', 0) for d in individual_stats_data.values())
                individual_rate = (total_awarded_ind / total_possible_ind * 100) if total_possible_ind > 0 else 0
                st.metric(label=f"{selected_user} の平均得点率", value=f"{individual_rate:.1f} %")

            # 「偏差値」を真ん中(col2)に配置
            with col2:
                user_total_score = sum(d.get('total_awarded', 0) for d in individual_stats_data.values())
                hensachi = calculate_hensachi(user_total_score, mean_score, std_dev_score)
                st.metric(label=f"{selected_user} の偏差値", value=f"{hensachi:.1f}")
        else:
            # 「全体」が選択されているときは、個人の欄は空けておく
            with col1:
                st.empty()
            with col2:
                st.empty()

        st.divider()
        if overall_stats_data:
            chart_image = create_performance_chart(individual_stats_data, overall_stats_data)
            if chart_image:
                st.image(chart_image)
        else:
            st.info("まだ解答データがなく、グラフを表示できません。")

        # --- フィードバック分析 ---
        st.divider()
        st.header("講義内容の評価・フィードバック")

        feedback_df = get_feedback_data()

        if feedback_df.empty:
            st.info("まだフィードバックがありません。")
        else:
            display_df = feedback_df[feedback_df['user_id'] == selected_user] if selected_user != "全体" else feedback_df

            if display_df.empty:
                st.warning(f"「{selected_user}」からのフィードバックはまだありません。")
            else:
                avg_rating = display_df['satisfaction_rating'].mean()
                metric_label = f"{selected_user} の平均満足度" if selected_user != "全体" else "全体の平均満足度"
                st.metric(label=metric_label, value=f"{avg_rating:.2f} / 5.0", delta=rating_to_stars(avg_rating))

                st.subheader("自由記述コメント一覧")
                for index, row in display_df.iterrows():
                    with st.expander(f"{row['submitted_at']} - {row['user_id']} (評価: {rating_to_stars(row['satisfaction_rating'])})"):
                        st.text(row['free_text_comment']) 
        
        st.divider()
        st.header("AIによる総合分析")
        
        if st.button("総合分析を実行する"):
            # 分析に必要なデータを準備
            target_scope = selected_user
            
            # 1. スコアデータ
            stats_data_for_ai = individual_stats_data if target_scope != "全体" else overall_stats_data
            score_df_for_ai = pd.DataFrame.from_dict(stats_data_for_ai, orient='index')
            if not score_df_for_ai.empty:
                score_df_for_ai['score_rate'] = (score_df_for_ai['total_awarded'] / score_df_for_ai['total_possible']) * 100

            # 2. フィードバックデータ
            all_feedback_df = get_feedback_data()
            feedback_df_for_ai = all_feedback_df[all_feedback_df['user_id'] == target_scope] if target_scope != "全体" else all_feedback_df

            # 3. 問題文データ
            conn = sqlite3.connect('qa.db');
            questions_df_for_ai = pd.read_sql_query("SELECT id, question_text FROM questions", conn).set_index('id')
            conn.close()

            # AI分析を実行
            with st.spinner("AIが成績とフィードバックを総合的に分析しています..."):
                analysis_result = generate_ai_analysis(target_scope, score_df_for_ai, feedback_df_for_ai, questions_df_for_ai)
                st.session_state.last_analysis_result = analysis_result
        
        # 分析結果を表示
        if 'last_analysis_result' in st.session_state:
            st.markdown(st.session_state.last_analysis_result)